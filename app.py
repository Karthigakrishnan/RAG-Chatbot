"""
DocMind RAG Chatbot
━━━━━━━━━━━━━━━━━━
No LangChain chains — uses FAISS + sentence-transformers + Groq SDK directly.
Zero import errors, fast indexing, multi-turn memory built-in.

✏️  Set your Groq API key below, then run:  streamlit run app.py
"""

import streamlit as st
import os, tempfile, time, json
from dotenv import load_dotenv
load_dotenv()  # loads .env file if present
import xml.etree.ElementTree as ET
from pathlib import Path

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  CONFIG  — paste your key here
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")  # set in .env
GROQ_MODEL   = "mixtral-8x7b-32768"     # or mixtral-8x7b-32768, llama3-8b-8192
TOP_K        = 4                     # chunks retrieved per query
CHUNK_SIZE   = 400                   # words per chunk
CHUNK_OVERLAP = 60                   # word overlap between chunks
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

st.set_page_config(
    page_title="DocMind · RAG",
    page_icon="🌿",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── CSS ───────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Serif+Display&family=DM+Sans:wght@300;400;500;600&family=JetBrains+Mono:wght@400;500&display=swap');

:root {
    --bg:       #f5f8f5;
    --surface:  #ffffff;
    --surf2:    #eef4ee;
    --surf3:    #e4ede4;
    --border:   #d0e4d0;
    --bord2:    #bcd4bc;
    --accent:   #2d6a4f;
    --acc-mid:  #40916c;
    --acc-lt:   #74c69d;
    --acc-xs:   #d8f3dc;
    --user-bg:  #e9f5ec;
    --text:     #192e19;
    --text2:    #3a5a3a;
    --muted:    #6b936b;
    --sans:     'DM Sans', sans-serif;
    --serif:    'DM Serif Display', serif;
    --mono:     'JetBrains Mono', monospace;
}

*, *::before, *::after { box-sizing: border-box; }
html, body, [class*="css"] {
    font-family: var(--sans) !important;
    background: var(--bg) !important;
    color: var(--text) !important;
}
.main, .stApp { background: var(--bg) !important; }
#MainMenu, footer { visibility: hidden; }
.block-container { padding: 2rem 2.5rem 5rem !important; max-width: 960px !important; }

/* sidebar */
[data-testid="stSidebar"] {
    background: var(--surface) !important;
    border-right: 1px solid var(--border) !important;
}
[data-testid="stSidebar"] > div:first-child { padding: 1.5rem 1.2rem !important; }

/* title */
@keyframes gentle-pulse { 0%,100%{opacity:1} 50%{opacity:.72} }
.app-title {
    font-family: var(--serif); font-size: 2rem;
    color: var(--accent); margin: 0 0 .15rem;
    animation: gentle-pulse 5s ease-in-out infinite;
}
.app-sub { font-size:.84rem; color:var(--muted); margin:0 0 .8rem; }

/* sidebar label */
.sb-label {
    font-family:var(--mono); font-size:.65rem;
    letter-spacing:.13em; text-transform:uppercase;
    color:var(--muted); margin:1rem 0 .35rem; display:block;
}

/* status pill */
.pill-green {
    display:inline-flex; align-items:center; gap:.38rem;
    background:var(--acc-xs); border:1px solid #b7dfc4;
    border-radius:999px; padding:.26rem .72rem;
    font-family:var(--mono); font-size:.7rem; color:var(--accent);
    margin-bottom:.6rem;
}
.pill-red {
    display:inline-flex; align-items:center; gap:.38rem;
    background:#fef2f2; border:1px solid #fca5a5;
    border-radius:999px; padding:.26rem .72rem;
    font-family:var(--mono); font-size:.7rem; color:#991b1b;
    margin-bottom:.6rem;
}
.sdot-green { width:7px; height:7px; background:var(--acc-lt); border-radius:50%; }
.sdot-red   { width:7px; height:7px; background:#f87171; border-radius:50%; }

/* stat chips */
.stat-row { display:flex; gap:.5rem; flex-wrap:wrap; margin:.35rem 0 .85rem; }
.stat-chip {
    background:var(--acc-xs); border:1px solid #b7dfc4;
    border-radius:999px; padding:.2rem .62rem;
    font-family:var(--mono); font-size:.68rem; color:var(--accent);
}

/* inputs */
input, textarea {
    background: var(--surf2) !important;
    border: 1px solid var(--bord2) !important;
    border-radius: 8px !important;
    color: var(--text) !important;
    font-family: var(--sans) !important;
    font-size: .9rem !important;
}
input:focus, textarea:focus {
    border-color: var(--acc-mid) !important;
    box-shadow: 0 0 0 3px rgba(64,145,108,.13) !important;
    outline: none !important;
}

/* selectbox */
.stSelectbox [data-baseweb="select"] > div {
    background: var(--surf2) !important;
    border: 1px solid var(--bord2) !important;
    border-radius: 8px !important;
    color: var(--text) !important;
}

/* buttons */
.stButton > button {
    background: var(--accent) !important;
    color: #fff !important; border: none !important;
    border-radius: 8px !important; font-family: var(--sans) !important;
    font-weight: 600 !important; font-size:.88rem !important;
    transition: all .17s !important; width:100%;
}
.stButton > button:hover {
    background: #235c42 !important;
    box-shadow: 0 3px 14px rgba(45,106,79,.22) !important;
    transform: translateY(-1px) !important;
}
.stButton > button:active { transform: translateY(0) !important; }
.ghost .stButton > button {
    background:transparent !important; color:var(--muted) !important;
    border:1px solid var(--bord2) !important;
}
.ghost .stButton > button:hover {
    color:#b84040 !important; border-color:#d47070 !important;
    box-shadow:none !important; transform:none !important;
}

/* progress */
.stProgress > div > div > div {
    background: linear-gradient(90deg,var(--accent),var(--acc-lt)) !important;
    border-radius: 999px !important;
}

/* file uploader shimmer */
@keyframes shimmer {
    0%,100%{ border-color:var(--bord2); box-shadow:none; }
    50%    { border-color:var(--acc-lt); box-shadow:0 0 16px rgba(116,198,157,.2); }
}
[data-testid="stFileUploader"] {
    background: var(--surface) !important;
    border: 2px dashed var(--bord2) !important;
    border-radius: 14px !important;
    animation: shimmer 3.5s ease-in-out infinite;
    padding: .6rem !important;
}

/* chat */
@keyframes fade-up {
    from{opacity:0;transform:translateY(8px)} to{opacity:1;transform:translateY(0)}
}
.chat-row { animation:fade-up .27s ease forwards; margin:.8rem 0; }
.user-lbl {
    font-family:var(--mono); font-size:.65rem; color:var(--acc-mid);
    letter-spacing:.1em; text-align:right; margin-bottom:.26rem;
}
.user-bub {
    background:var(--user-bg); border:1px solid #bcddc8;
    border-radius:14px 14px 3px 14px;
    padding:.8rem 1rem; max-width:76%; margin-left:auto;
    font-size:.9rem; line-height:1.65; color:var(--text);
}
.bot-lbl {
    font-family:var(--mono); font-size:.65rem;
    color:var(--muted); letter-spacing:.1em; margin-bottom:.26rem;
}
.bot-bub {
    background:var(--surface); border:1px solid var(--border);
    border-left:3px solid var(--acc-lt);
    border-radius:3px 14px 14px 14px;
    padding:.8rem 1rem; max-width:85%;
    font-size:.9rem; line-height:1.72; color:var(--text);
    box-shadow:0 1px 5px rgba(0,0,0,.04);
    white-space: pre-wrap;
}

/* typing dots */
@keyframes dot-up {
    0%,80%,100%{transform:translateY(0);opacity:.3}
    40%{transform:translateY(-5px);opacity:1}
}
.typing-wrap {
    display:flex; align-items:center; gap:5px;
    background:var(--surface); border:1px solid var(--border);
    border-left:3px solid var(--acc-lt);
    border-radius:3px 14px 14px 14px;
    padding:.68rem 1rem; width:fit-content;
}
.tdot {
    width:7px; height:7px; background:var(--acc-lt);
    border-radius:50%; animation:dot-up 1.2s infinite ease-in-out;
}
.tdot:nth-child(2){animation-delay:.18s}
.tdot:nth-child(3){animation-delay:.36s}

/* source card */
.src-card {
    background:var(--surf2); border:1px solid var(--border);
    border-left:3px solid var(--bord2); border-radius:6px;
    padding:.62rem .88rem; margin-bottom:.42rem;
    font-family:var(--mono); font-size:.77rem;
    color:var(--text2); line-height:1.6;
}
.src-meta {
    font-size:.62rem; color:var(--accent);
    letter-spacing:.08em; display:block; margin-bottom:.26rem;
}

/* format pills */
.fmt-pills { display:flex; flex-wrap:wrap; gap:.32rem; justify-content:center; margin-top:.85rem; }
.fmt-pill {
    background:var(--surf3); border:1px solid var(--border);
    border-radius:999px; padding:.15rem .55rem;
    font-family:var(--mono); font-size:.66rem; color:var(--text2);
}

hr { border:none; border-top:1px solid var(--border) !important; margin:1rem 0 !important; }
</style>
""", unsafe_allow_html=True)


# ── Session state ─────────────────────────────────────────────────────────────
for k, v in {
    "chat_history":  [],      # [{role, content, sources}]
    "chunks":        [],      # list of str
    "embeddings_mat": None,   # numpy array
    "index":         None,    # faiss index
    "processed":     False,
    "doc_stats":     {},
    "model_name":    GROQ_MODEL,
}.items():
    if k not in st.session_state:
        st.session_state[k] = v


# ── Cached embedding model ────────────────────────────────────────────────────
@st.cache_resource(show_spinner=False)
def load_embed_model():
    from sentence_transformers import SentenceTransformer
    return SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")


# ── Chunking ──────────────────────────────────────────────────────────────────
def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunk = " ".join(words[i: i + size])
        if chunk.strip():
            chunks.append(chunk)
        i += size - overlap
    return chunks


# ── Text extraction ───────────────────────────────────────────────────────────
def extract_text(path: str, name: str) -> str:
    ext = Path(name).suffix.lower()
    try:
        if ext == ".pdf":
            try:
                import fitz
                doc = fitz.open(path)
                return "\n\n".join(
                    f"[Page {i+1}]\n{p.get_text()}"
                    for i, p in enumerate(doc) if p.get_text().strip()
                )
            except ImportError:
                import pdfplumber
                with pdfplumber.open(path) as pdf:
                    return "\n\n".join(
                        f"[Page {i+1}]\n{p.extract_text()}"
                        for i, p in enumerate(pdf.pages) if p.extract_text()
                    )
        elif ext == ".docx":
            from docx import Document
            doc = Document(path)
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for t in doc.tables:
                for row in t.rows:
                    r = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if r: parts.append(r)
            return "\n\n".join(parts)
        elif ext in [".pptx", ".ppt"]:
            from pptx import Presentation
            prs = Presentation(path)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                parts = [f"[Slide {i}]"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t: parts.append(t)
                if len(parts) > 1:
                    slides.append("\n".join(parts))
            return "\n\n".join(slides)
        elif ext in [".xlsx", ".xls"]:
            import pandas as pd
            xf = pd.ExcelFile(path)
            return "\n\n".join(
                f"[Sheet: {n}]\n{xf.parse(n).to_string(index=False)}"
                for n in xf.sheet_names
            )
        elif ext == ".csv":
            import pandas as pd
            return pd.read_csv(path).to_string(index=False)
        elif ext == ".json":
            with open(path, encoding="utf-8", errors="ignore") as f:
                return json.dumps(json.load(f), indent=2)
        elif ext == ".xml":
            root = ET.parse(path).getroot()
            return "\n".join(
                f"{e.tag}: {e.text.strip()}"
                for e in root.iter() if e.text and e.text.strip()
            )
        elif ext in [".html", ".htm"]:
            try:
                from bs4 import BeautifulSoup
                with open(path, encoding="utf-8", errors="ignore") as f:
                    return BeautifulSoup(f.read(), "html.parser").get_text("\n")
            except ImportError:
                import re, html as _h
                with open(path, encoding="utf-8", errors="ignore") as f:
                    return _h.unescape(re.sub(r"<[^>]+>", " ", f.read()))
        else:
            with open(path, encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception as e:
        st.warning(f"⚠️ Could not parse **{name}**: {e}")
        return ""


# ── Build FAISS index ─────────────────────────────────────────────────────────
def build_index(all_texts: list[str]):
    import faiss, numpy as np

    model = load_embed_model()

    # Chunk all documents
    all_chunks = []
    for text in all_texts:
        all_chunks.extend(chunk_text(text))

    if not all_chunks:
        raise ValueError("No text chunks could be created.")

    # Embed (batched — fast)
    vecs = model.encode(
        all_chunks,
        batch_size=64,
        show_progress_bar=False,
        convert_to_numpy=True,
        normalize_embeddings=True,
    ).astype("float32")

    # FAISS inner-product index (= cosine sim on normalised vecs)
    dim = vecs.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(vecs)

    return all_chunks, vecs, index


# ── Retrieve top-k chunks ─────────────────────────────────────────────────────
def retrieve(query: str, top_k: int = TOP_K) -> list[str]:
    import numpy as np
    model = load_embed_model()
    q_vec = model.encode(
        [query], normalize_embeddings=True, convert_to_numpy=True
    ).astype("float32")
    scores, idxs = st.session_state.index.search(q_vec, top_k)
    chunks = st.session_state.chunks
    return [chunks[i] for i in idxs[0] if i < len(chunks)]


# ── Groq answer ───────────────────────────────────────────────────────────────
def ask_groq(question: str, context_chunks: list[str]) -> str:
    from groq import Groq

    client = Groq(api_key=GROQ_API_KEY)
    context = "\n\n---\n\n".join(
        f"[Chunk {i+1}]\n{c}" for i, c in enumerate(context_chunks)
    )

    # Build message history for multi-turn memory
    messages = [
        {
            "role": "system",
            "content": (
                "You are a precise document Q&A assistant. "
                "Answer ONLY using the provided document context. "
                "If the answer is not in the context, say: "
                "'The document does not contain information about this.' "
                "Be concise and accurate."
                f"\n\nDOCUMENT CONTEXT:\n{context}"
            ),
        }
    ]

    # Add previous turns (last 6 exchanges = 12 messages)
    for msg in st.session_state.chat_history[-12:]:
        messages.append({
            "role": msg["role"] if msg["role"] in ("user", "assistant") else "user",
            "content": msg["content"],
        })

    # Add current question
    messages.append({"role": "user", "content": question})

    resp = client.chat.completions.create(
        model=st.session_state.model_name,
        messages=messages,
        temperature=0.2,
        max_tokens=1024,
    )
    return resp.choices[0].message.content.strip()


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  SIDEBAR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
with st.sidebar:
    st.markdown("### 🌿 DocMind")
    st.markdown(
        '<p style="font-size:.77rem;color:#6b936b;margin-top:-.2rem;margin-bottom:.4rem">'
        "RAG · Groq · FAISS</p>",
        unsafe_allow_html=True,
    )
    st.markdown("---")

    # API key status
    st.markdown('<span class="sb-label">API Status</span>', unsafe_allow_html=True)
    key_ok = bool(GROQ_API_KEY and GROQ_API_KEY != "gsk_your_key_here")
    if key_ok:
        st.markdown(
            '<div class="pill-green"><div class="sdot-green"></div>Groq connected</div>',
            unsafe_allow_html=True,
        )
    else:
        st.markdown(
            '<div class="pill-red"><div class="sdot-red"></div>No API key</div>',
            unsafe_allow_html=True,
        )
        st.caption("Open app.py and set GROQ_API_KEY on line 22.")

    # Model selector
    st.markdown('<span class="sb-label">Model</span>', unsafe_allow_html=True)
    model_options = [
        "llama3-70b-8192",
        "llama3-8b-8192",
        "mixtral-8x7b-32768",
        "gemma2-9b-it",
        "llama-3.1-70b-versatile",
        "llama-3.1-8b-instant",
    ]
    chosen = st.selectbox(
        "mdl", model_options,
        index=model_options.index(st.session_state.model_name)
              if st.session_state.model_name in model_options else 0,
        label_visibility="collapsed",
    )
    st.session_state.model_name = chosen

    st.markdown("---")

    # Doc stats
    if st.session_state.processed:
        s = st.session_state.doc_stats
        st.markdown(
            f'<div class="stat-row">'
            f'<span class="stat-chip">📄 {s.get("files",0)} files</span>'
            f'<span class="stat-chip">🧩 {s.get("chunks",0)} chunks</span>'
            f'</div>'
            f'<div class="stat-row">'
            f'<span class="stat-chip">📝 {s.get("words",0):,} words</span>'
            f'</div>',
            unsafe_allow_html=True,
        )
        st.markdown("---")

    if st.session_state.chat_history:
        st.markdown('<div class="ghost">', unsafe_allow_html=True)
        if st.button("🗑  Clear Chat", use_container_width=True):
            st.session_state.chat_history = []
            st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)

    if st.session_state.processed:
        st.markdown('<div class="ghost">', unsafe_allow_html=True)
        if st.button("↩  New Document", use_container_width=True):
            st.session_state.update({
                "chunks": [], "embeddings_mat": None, "index": None,
                "processed": False, "doc_stats": {}, "chat_history": [],
            })
            st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("---")
    st.markdown(
        '<p style="font-size:.68rem;color:#afc8af;text-align:center;line-height:1.9">'
        "PDF · DOCX · PPTX · XLSX<br>CSV · TXT · MD · JSON · XML · HTML</p>",
        unsafe_allow_html=True,
    )


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  MAIN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
st.markdown('<h1 class="app-title">🌿 DocMind</h1>', unsafe_allow_html=True)
st.markdown(
    '<p class="app-sub">Ask questions grounded in your documents — powered by Groq</p>',
    unsafe_allow_html=True,
)
st.markdown("---")

# ━━ UPLOAD ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
if not st.session_state.processed:

    _, col, _ = st.columns([1, 2.6, 1])
    with col:
        st.markdown("""
        <div style="text-align:center;padding:1.8rem 0 .8rem">
            <div style="font-size:2.5rem;margin-bottom:.5rem">📂</div>
            <div style="font-family:'DM Serif Display',serif;font-size:1.45rem;
                        color:#2d6a4f;margin-bottom:.4rem">Upload your documents</div>
            <div style="font-size:.83rem;color:#6b936b;line-height:1.7;
                        max-width:330px;margin:0 auto 1.4rem">
                Drop files below and click
                <strong style="color:#2d6a4f">Process</strong>.
                Answers are grounded exclusively in your document content.
            </div>
        </div>
        """, unsafe_allow_html=True)

        uploaded = st.file_uploader(
            "upload", accept_multiple_files=True,
            type=["pdf","docx","pptx","ppt","xlsx","xls",
                  "csv","txt","md","json","xml","html","htm"],
            label_visibility="collapsed",
        )

        st.markdown("""
        <div class="fmt-pills">
            <span class="fmt-pill">PDF</span><span class="fmt-pill">DOCX</span>
            <span class="fmt-pill">PPTX</span><span class="fmt-pill">XLSX</span>
            <span class="fmt-pill">CSV</span><span class="fmt-pill">TXT</span>
            <span class="fmt-pill">MD</span><span class="fmt-pill">JSON</span>
            <span class="fmt-pill">XML</span><span class="fmt-pill">HTML</span>
        </div>
        """, unsafe_allow_html=True)

        st.markdown("<br>", unsafe_allow_html=True)

        if uploaded:
            if not key_ok:
                st.error("❌ Set your GROQ_API_KEY in app.py line 22 first.")
            else:
                if st.button("⚡  Process Documents", use_container_width=True):
                    prog = st.progress(0, text="Extracting text…")
                    all_texts = []

                    for i, uf in enumerate(uploaded):
                        prog.progress(
                            int(i / len(uploaded) * 40),
                            text=f"Reading {uf.name}…",
                        )
                        with tempfile.NamedTemporaryFile(
                            delete=False, suffix=Path(uf.name).suffix
                        ) as tmp:
                            tmp.write(uf.getvalue())
                            tmp_path = tmp.name
                        txt = extract_text(tmp_path, uf.name)
                        os.unlink(tmp_path)
                        if txt.strip():
                            all_texts.append(txt)
                        else:
                            st.warning(f"No text found in {uf.name}")

                    if all_texts:
                        prog.progress(50, text="Loading embedding model…")
                        try:
                            chunks, vecs, index = build_index(all_texts)
                            prog.progress(92, text="Finalising…")
                            st.session_state.update({
                                "chunks":        chunks,
                                "embeddings_mat": vecs,
                                "index":         index,
                                "processed":     True,
                                "chat_history":  [],
                                "doc_stats": {
                                    "files":  len(uploaded),
                                    "chunks": len(chunks),
                                    "words":  sum(len(t.split()) for t in all_texts),
                                },
                            })
                            prog.progress(100, text="Done ✓")
                            time.sleep(0.3)
                            prog.empty()
                            st.rerun()
                        except Exception as e:
                            prog.empty()
                            st.error(f"❌ {e}")
                    else:
                        prog.empty()
                        st.error("No text could be extracted from the uploaded files.")

# ━━ CHAT ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
else:
    model_label = st.session_state.model_name

    # Render history
    for i, msg in enumerate(st.session_state.chat_history):
        if msg["role"] == "user":
            st.markdown(f"""
            <div class="chat-row">
                <div class="user-lbl">YOU</div>
                <div class="user-bub">{msg["content"]}</div>
            </div>""", unsafe_allow_html=True)
        else:
            st.markdown(f"""
            <div class="chat-row">
                <div class="bot-lbl">🌿 DOCMIND · {model_label}</div>
                <div class="bot-bub">{msg["content"]}</div>
            </div>""", unsafe_allow_html=True)

            if msg.get("sources"):
                n = sum(
                    1 for m in st.session_state.chat_history[:i+1]
                    if m["role"] == "assistant"
                )
                with st.expander(
                    f"📎 {len(msg['sources'])} source chunks — answer #{n}",
                    expanded=False,
                ):
                    for j, src in enumerate(msg["sources"]):
                        snippet = src[:380] + ("…" if len(src) > 380 else "")
                        st.markdown(f"""
                        <div class="src-card">
                            <span class="src-meta">CHUNK {j+1}</span>{snippet}
                        </div>""", unsafe_allow_html=True)

    # Suggested starters
    if not st.session_state.chat_history:
        st.markdown(
            '<p style="font-family:\'JetBrains Mono\',monospace;font-size:.68rem;'
            'letter-spacing:.12em;color:#afc8af;margin:1rem 0 .5rem">'
            'SUGGESTED QUESTIONS</p>',
            unsafe_allow_html=True,
        )
        q_cols = st.columns(2)
        for idx, q in enumerate([
            "Summarize the main topics.",
            "What are the key findings?",
            "List the important concepts.",
            "What conclusions are drawn?",
        ]):
            with q_cols[idx % 2]:
                if st.button(q, key=f"sq_{idx}", use_container_width=True):
                    st.session_state.chat_history.append(
                        {"role": "user", "content": q, "sources": []}
                    )
                    ph = st.empty()
                    ph.markdown(
                        '<div class="typing-wrap">'
                        '<div class="tdot"></div><div class="tdot"></div>'
                        '<div class="tdot"></div></div>',
                        unsafe_allow_html=True,
                    )
                    try:
                        srcs = retrieve(q)
                        ans  = ask_groq(q, srcs)
                    except Exception as e:
                        ans, srcs = f"❌ {e}", []
                    ph.empty()
                    st.session_state.chat_history.append(
                        {"role": "assistant", "content": ans, "sources": srcs}
                    )
                    st.rerun()

    # Input bar
    st.markdown("---")
    c1, c2 = st.columns([6, 1])
    with c1:
        user_q = st.text_input(
            "q",
            placeholder="Ask anything about your document…",
            label_visibility="collapsed",
            key="chat_input",
        )
    with c2:
        send = st.button("Send", use_container_width=True)

    if send and user_q.strip():
        q = user_q.strip()
        st.session_state.chat_history.append(
            {"role": "user", "content": q, "sources": []}
        )
        ph = st.empty()
        ph.markdown(
            '<div class="typing-wrap">'
            '<div class="tdot"></div><div class="tdot"></div>'
            '<div class="tdot"></div></div>',
            unsafe_allow_html=True,
        )
        try:
            srcs = retrieve(q)
            ans  = ask_groq(q, srcs)
        except Exception as e:
            err = str(e)
            if any(x in err.lower() for x in ["api_key", "401", "auth", "invalid", "forbidden"]):
                ans = "❌ Invalid Groq API key. Update GROQ_API_KEY in app.py"
            elif "rate" in err.lower():
                ans = "⏳ Rate limit hit — wait a moment and retry."
            elif "model" in err.lower():
                ans = f"❌ Model error: {err}"
            else:
                ans = f"❌ {err}"
            srcs = []
        ph.empty()
        st.session_state.chat_history.append(
            {"role": "assistant", "content": ans, "sources": srcs}
        )
        st.rerun()
    elif send:
        st.warning("Please type a question first.")